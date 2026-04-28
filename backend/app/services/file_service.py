# app/services/file_service.py
import os
import logging
import io
from fastapi import UploadFile, HTTPException
from typing import Tuple
import uuid
from app.config.logging_config import setup_logging
from app.services.supabase_connection import SupabaseConnection

# Text extraction imports
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation

# Setup logging
setup_logging() 
logger = logging.getLogger(__name__)


class FileService:
    def __init__(self):
        """Initialize FileService instance."""
        self.supabase_connection = SupabaseConnection()
        self.supabase = self.supabase_connection.client
        self.bucket_name = self.supabase_connection.get_bucket_name()
    
    def get_file_type(self, filename: str) -> str:
        """Extract file type from filename."""
        logger.debug(f"Extracting file type from filename: {filename}")
        extension = os.path.splitext(filename)[1].lower()
        
        if extension == ".pdf":
            logger.debug(f"File type detected: pdf")
            return "pdf"
        elif extension == ".pptx":
            logger.debug(f"File type detected: pptx")
            return "pptx"
        elif extension == ".docx":
            logger.debug(f"File type detected: docx")
            return "docx"
        elif extension == ".txt":
            logger.debug(f"File type detected: txt")
            return "txt"
        else:
            logger.warning(f"Unsupported file type: {extension}")
            raise HTTPException(
                status_code=400, 
                detail="Unsupported file type. Supported types: pdf, pptx, docx, txt"
            )
    
    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """
        Extract text from PDF file content.
        
        Args:
            file_content (bytes): PDF file content
            
        Returns:
            str: Extracted text content
        """
        logger.info("Extracting text from PDF file")
        try:
            # Try with pdfplumber first (better for complex layouts)
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                text_content = []
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
                
                if text_content:
                    extracted_text = "\n\n".join(text_content)
                    logger.info(f"Successfully extracted {len(extracted_text)} characters using pdfplumber")
                    return extracted_text
            
            # Fallback to PyPDF2 if pdfplumber fails
            logger.info("Falling back to PyPDF2 for PDF extraction")
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text_content = []
            
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
            
            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters using PyPDF2")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to extract text from PDF: {str(e)}"
            )
    
    def extract_text_from_docx(self, file_content: bytes) -> str:
        """
        Extract text from DOCX file content.
        
        Args:
            file_content (bytes): DOCX file content
            
        Returns:
            str: Extracted text content
        """
        logger.info("Extracting text from DOCX file")
        try:
            doc = Document(io.BytesIO(file_content))
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters from DOCX")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to extract text from DOCX: {str(e)}"
            )
    
    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """
        Extract text from PPTX file content.
        
        Args:
            file_content (bytes): PPTX file content
            
        Returns:
            str: Extracted text content
        """
        logger.info("Extracting text from PPTX file")
        try:
            prs = Presentation(io.BytesIO(file_content))
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters from PPTX")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to extract text from PPTX: {str(e)}"
            )
    
    def extract_text_from_txt(self, file_content: bytes) -> str:
        """
        Extract text from TXT file content.
        
        Args:
            file_content (bytes): TXT file content
            
        Returns:
            str: Extracted text content
        """
        logger.info("Extracting text from TXT file")
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    extracted_text = file_content.decode(encoding)
                    logger.info(f"Successfully extracted {len(extracted_text)} characters from TXT using {encoding} encoding")
                    return extracted_text
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use utf-8 with error handling
            extracted_text = file_content.decode('utf-8', errors='replace')
            logger.warning("Used UTF-8 with error replacement for TXT file")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from TXT: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to extract text from TXT: {str(e)}"
            )
    
    def extract_text_content(self, file_content: bytes, file_type: str) -> str:
        """
        Extract text content from file based on file type.
        
        Args:
            file_content (bytes): File content as bytes
            file_type (str): Type of file (pdf, docx, pptx, txt)
            
        Returns:
            str: Extracted and formatted text content
        """
        logger.info(f"Extracting text content for file type: {file_type}")
        
        if file_type == "pdf":
            return self.extract_text_from_pdf(file_content)
        elif file_type == "docx":
            return self.extract_text_from_docx(file_content)
        elif file_type == "pptx":
            return self.extract_text_from_pptx(file_content)
        elif file_type == "txt":
            return self.extract_text_from_txt(file_content)
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type for text extraction: {file_type}"
            )

    async def save_upload_file(self, file: UploadFile) -> Tuple[str, str]:
        """
        Save uploaded file to Supabase storage and return file path and type.
        
        Returns:
            Tuple[str, str]: (file_path, file_type)
        """
        logger.info(f"Processing file upload: {file.filename}")
        file_type = self.get_file_type(file.filename)
        
        # Generate unique filename
        unique_filename = f"{uuid.uuid4()}.{file_type}"
        logger.debug(f"Generated unique filename: {unique_filename}")
        
        # Read file content
        file_content = await file.read()
        logger.debug(f"Read file content, size: {len(file_content)} bytes")
        
        # Upload to Supabase Storage
        try:
            logger.info(f"Uploading file to Supabase bucket: {self.bucket_name}")
            result = self.supabase.storage.from_(self.bucket_name).upload(
                unique_filename,
                file_content,
                {"content-type": file.content_type}
            )
            
            # Get the public URL
            file_path = self.supabase.storage.from_(self.bucket_name).get_public_url(unique_filename)
            logger.info(f"File uploaded successfully. Path: {file_path}")
            
            return file_path, file_type
            
        except Exception as e:
            logger.error(f"Failed to upload file to Supabase: {str(e)}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to upload file to Supabase: {str(e)}"
            )
    
    async def extract_text_from_upload(self, file: UploadFile) -> Tuple[str, str]:
        """
        Extract text content from uploaded file without saving to storage.
        
        Args:
            file (UploadFile): The uploaded file
            
        Returns:
            Tuple[str, str]: (extracted_text, file_type)
        """
        logger.info(f"Extracting text from uploaded file: {file.filename}")
        file_type = self.get_file_type(file.filename)
        
        # Read file content
        file_content = await file.read()
        logger.debug(f"Read file content, size: {len(file_content)} bytes")
        
        # Extract text content
        extracted_text = self.extract_text_content(file_content, file_type)
        
        # Reset file position for potential future reads
        await file.seek(0)
        
        return extracted_text, file_type